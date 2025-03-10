import datetime
import os
import asyncio
import threading
from openai import OpenAI
from pptx import Presentation
from .file_handler import FileHandler
from utils.logger import get_logger
from utils.helpers import create_file, extract_tables, get_resource_path
from config import ConfigManager

logger = get_logger()

class LLMHandler:
    def __init__(self, config_manager: ConfigManager):
        logger.debug("Initializing LLMHandler")
        self.config_manager = config_manager
        
        # Initialize API configuration
        api_config = {
            'base_url': self.config_manager.get('api_config', 'base_url'),
            'api_key': self.config_manager.get('api_config', 'api_key'),
            'model': self.config_manager.get('api_config', 'model'),
            'temperature': self.config_manager.get('api_config', 'temperature'),
            'top_p': self.config_manager.get('api_config', 'top_p'),
            'max_tokens': self.config_manager.get('api_config', 'max_tokens'),
            'stream': self.config_manager.get('api_config', 'stream')
        }
        
        # Store API parameters
        for key, value in api_config.items():
            setattr(self, key, value)
            
        # Initialize OpenAI client
        self.client = OpenAI(base_url=self.base_url, api_key=self.api_key)
        
        # Log configuration (excluding sensitive data)
        logger.debug(f"Model: {self.model} Temperature: {self.temperature} Top P: {self.top_p} Max Tokens: {self.max_tokens} Stream: {self.stream}")

    async def _call_api(self, prompt, messages=None,role="user"):
        if messages is None:
            messages = [{"role": role, "content": prompt}]
            
        try:
            completion = await asyncio.to_thread(
                self.client.chat.completions.create,
                model=self.model,
                messages=messages,
                temperature=self.temperature,
                top_p=self.top_p,
                max_tokens=self.max_tokens,
                stream=self.stream
            )
            
            response = ""
            for chunk in completion:
                if chunk.choices[0].delta.content is not None:
                    response += chunk.choices[0].delta.content

            return response
        except Exception as e:
            logger.error(f"Error calling API: {str(e)}")
            raise

    async def fetch_intent_prompt(self):
        intent_prompt_path = get_resource_path("resources/prompts/intent_selector.json")
        
        try:
            async with asyncio.Semaphore(1):  # Prevents multiple concurrent file accesses
                with open(intent_prompt_path, 'r') as file:
                    content = await asyncio.to_thread(file.read)
                    return content
        except Exception as e:
            logger.error(f"Error reading intent prompt: {str(e)}")
            raise

    async def init_prompt(self):
        logger.debug("Initializing Intent for LLMHandler")
        
        try:
            intent_prompt = await self.fetch_intent_prompt()
            response = await self._call_api(intent_prompt)
            return response
        except Exception as e:
            logger.error(f"Error initializing prompt: {str(e)}")
            raise

    def run_init_prompt(self):
        loop = asyncio.new_event_loop()
        threading.Thread(target=lambda: loop.run_until_complete(self.init_prompt()), daemon=True).start()

    def response(self, text):
        logger.info(f"Calling LLM endpoint {self.base_url}")
        logger.debug(f"Requesting completion for text: {text}")
        
        # Create synchronous version for backward compatibility
        completion = self.client.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": text}],
            temperature=self.temperature,
            top_p=self.top_p,
            max_tokens=self.max_tokens,
            stream=self.stream
        )
        
        response = ""
        for chunk in completion:
            if chunk.choices[0].delta.content is not None:
                response += chunk.choices[0].delta.content
                
        logger.debug(f"Received response: {response}")
        return response

    def create_listing(self, listing_type):
        logger.info(f"Creating {listing_type} listing")
        
        prompt = """
            Create a request listing for a {} data engineer. The list must have 
            'must' and 'should' criteria. Make the listing 500 characters or less.
        """.format(listing_type)
        
        return self.response(prompt)

    def review_cv(self, cv_text, listing):
        logger.info("Reviewing CV against listing")
        
        # Handle both string and list input for cv_text
        if isinstance(cv_text, list):
            cv_text = "\n".join(cv_text)
            
        prompt = """
            Given the following CV: 
            ---{}---
            and the following requirements listing: 
            ---{}---
            determine if the CV owner would be accepted to the listing or not. The output follows the following structure: 
            1. Name the candidate and the title of the listing role
            2. Use 'Accept' or 'Deny' as your answer
            3. Explain your decision with less than 100 characters.
        """.format(cv_text, listing)
        
        return self.response(prompt)
    
    def process_cv_batch(self, cv_files_directory, listings=None):
        logger.info(f"Processing CV batch from {cv_files_directory}")
        results = {}
        
        # Get CV files
        cv_files = self._get_cv_files(cv_files_directory)
        if not cv_files:
            logger.warning("No valid CV files found in directory")
            raise ValueError("No valid CV files found in directory")
        
        # Create listings if not provided
        if not listings:
            listings = self._create_default_listings(cv_files_directory)
        
        # Process each CV
        for file_name, file_path in cv_files.items():
            logger.debug(f"Processing CV file: {file_name}")
            cv_text = self.extract_text_from_pptx(file_path)
            results[file_name] = {}
            
            for listing_name, listing_text in listings.items():
                logger.debug(f"Reviewing {file_name} against {listing_name} listing")
                results[file_name][listing_name] = self.review_cv(cv_text, listing_text)
        
        return results, listings
    
    def _get_cv_files(self, directory):
        files = os.listdir(directory)
        cv_files = {}
        
        for file in files:
            if file.endswith('.pptx') and not file.startswith('corrected_'):
                file_path = os.path.join(directory, file)
                cv_files[file] = file_path
                
        return cv_files
    
    def _create_default_listings(self, directory):
        logger.info("Creating default listings")
        
        listings = {
            'generic': self.create_listing('generic'),
            'highly_experienced': self.create_listing('highly experienced (senior)'),
        }
        
        # Try to extract from PDF if available
        files = os.listdir(directory)
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        
        if pdf_files:
            try:
                pdf_tables = extract_tables(os.path.join(directory, pdf_files[0]))
                if pdf_tables:
                    listings['pdf_based'] = self.table_analysis(pdf_tables)
            except Exception as e:
                logger.error(f"Error processing PDF: {str(e)}")
                
        return listings
    
    def extract_text_from_pptx(self, pptx_path: str) -> str:
        """Extract and correct text from a PowerPoint file"""
        logger.info(f"Extracting text from {pptx_path}")
        
        try:
            prs = Presentation(pptx_path)
            all_text = []
            
            for slide_num, slide in enumerate(prs.slides):
                logger.debug(f"Processing slide {slide_num+1}")
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame is not None:
                        text = shape.text_frame.text.strip()
                        if text:
                            correct_text = self.spelling_and_grammar_check(text)
                            slide_text.append(correct_text)
                            
                if slide_text:
                    all_text.append("\n".join(slide_text))

            return "\n\n".join(all_text)
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {str(e)}")
            raise
    
    def spelling_and_grammar_check(self, text: str):
        logger.debug(f"Checking spelling and grammar")
        
        prompt = """
            Correct the spelling and grammar of the following text from a CV. Return ONLY the corrected text, no explanation is needed.
            If the original text is already correct or empty, return the same content. Do not add any additional text including explanation.
            Text: "{}"
            Example:
            Original Text: "Kandidat One"
            Return: 'Kandidat One'          
        """.format(text)
        
        return self.response(prompt)

    def table_analysis(self, tables: list) -> str:
        logger.info(f"Analyzing {len(tables)} tables")
        
        # Convert table objects to string representation
        table_strings = []
        for i, table in enumerate(tables):
            table_strings.append(f"Table {i+1}:\n{table.to_string()}")
        
        prompt = """
            Given the following tables:
            {}
            Analyze the tables and determine the 'must' and 'should' criteria for the requirements listing. 
            The list must have 'must' and 'should' criteria. Make the listing 500 characters or less.
        """.format("\n\n".join(table_strings))
        
        return self.response(prompt)

    async def _process_message_intent(self, user_message,file_handler: FileHandler):
        user_message_lower = user_message.lower()
        
        if "create listing" in user_message_lower or "job listing" in user_message_lower:
            # Look for listing type in message
            if "senior" in user_message_lower or "experienced" in user_message_lower:
                return self.create_listing("highly experienced (senior)")
            else:
                return self.create_listing("generic")
                
        elif "review" in user_message_lower or "resume" in user_message_lower:
            if not file_handler.get_uploaded_file_path():
                return "Please upload a CV file first."
            
            # Extract CV text
            try:
                cv_text = self.extract_text_from_pptx(file_handler.get_uploaded_file_path())
                listing = self.create_listing("generic")  # Create a generic listing
                result = self.review_cv(cv_text, listing)
                file_handler.reset_uploaded_file_path()
                return result
            except Exception as e:
                logger.error(f"Error processing CV: {str(e)}")
                return f"Error processing the CV: {str(e)}"
                
        elif "process cv batch" in user_message_lower or "batch process" in user_message_lower:
            # Get directory from the uploaded file's directory
            if not file_handler.get_uploaded_file_path():
                return "Please upload at least one CV file first."
                
            directory = os.path.dirname(file_handler.get_uploaded_file_path())
            try:
                results, listings = self.process_cv_batch(directory)
                
                # Save results to download directory
                result_text = "Batch Processing Results:\n\n"
                for cv_id, cv_results in results.items():
                    result_text += f"## {cv_id}\n"
                    for listing_name, review in cv_results.items():
                        result_text += f"- {listing_name}: {review}\n"
                    result_text += "\n"
                
                result_file = os.path.join(file_handler.storage_directory, f"batch_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt")
                with open(result_file, "w") as f:
                    f.write(result_text)
                
                self.downloaded_file_path = result_file
                return f"Batch processing completed. Results saved to {os.path.basename(result_file)}"
            except Exception as e:
                logger.error(f"Error in batch processing: {str(e)}")
                return f"Error in batch processing: {str(e)}"
                
        elif any(keyword in user_message_lower for keyword in ["spell","spelling", "grammar", "correction","correct","spellings","mistakes"]):
            if not file_handler.get_uploaded_file_path():
                return "Please upload a CV file first."
            try:
                cv_text = self.extract_text_from_pptx(file_handler.get_uploaded_file_path())
                result = self.spelling_and_grammar_check(cv_text)
                created = create_file(f"outputs", result)
                if created:
                    file_handler.download_file(f"outputs")
                file_handler.reset_uploaded_file_path()
                return result
            except Exception as e:
                logger.error(f"Error processing CV: {str(e)}")
                return f"Error processing the CV: {str(e)}"
            
        elif "table analysis" in user_message_lower or "analyze table" in user_message_lower:
            if not file_handler.get_uploaded_file_path():
                return "Please upload a PDF file containing tables first."
                
            try:
                tables = extract_tables(file_handler.get_uploaded_file_path())
                if not tables:
                    return "No tables found in the uploaded PDF."
                analysis = self.table_analysis(tables)
                return analysis
            except Exception as e:
                logger.error(f"Error analyzing tables: {str(e)}")
                return f"Error analyzing tables: {str(e)}"
        
        return self.response(user_message)